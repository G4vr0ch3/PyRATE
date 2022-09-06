#!/usr/bin/python3.10


################################################################################
#                                                                              #
#                                                                              #
#                   GNU AFFERO GENERAL PUBLIC LICENSE                          #
#                       Version 3, 19 November 2007                            #
#                                                                              #
#    This programms aims at sanitizing .odt and .ott documents.                #
#    Copyright (C) 2022 Gavroche, Roxane                                       #
#                                                                              #
#    This program is free software: you can redistribute it and/or modify      #
#    it under the terms of the GNU Affero General Public License as published  #
#    by the Free Software Foundation, either version 3 of the License, or      #
#    (at your option) any later version.                                       #
#                                                                              #
#     This program is distributed in the hope that it will be useful,          #
#    but WITHOUT ANY WARRANTY; without even the implied warranty of            #
#    MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the             #
#    GNU Affero General Public License for more details.                       #
#                                                                              #
#    You should have received a copy of the GNU Affero General Public License  #
#    along with this program.  If not, see <https://www.gnu.org/licenses/>.    #
#                                                                              #
#                                                                              #
################################################################################


from os import remove
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE


################################################################################


def get_layout(path):

    layout = []
    img_id = 0

    # Open source document
    pres = Presentation(path)

    # Get slides
    slides = pres.slides

    for slide in slides:

        s_ctnt = {"slide": slides.index(slide),
                "shapes": []
                }

        for shape in slide.shapes:

            shp_ctnt = {
                    "id": slide.shapes.index(shape),
                    "ctnt": {},
                    }

            # Slide's texts
            if shape.has_text_frame:
                txt = []
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        txt.append(run.text)

                shp_ctnt["ctnt"] = {"type": 'text', "ctnt": ''.join(txt)}

            # Slide's tables
            if shape.has_table:

                table = shape.table
                dim = len(table.rows), len(table.columns)

                cells = [(i, j, table.cell(i,j).text) for i in range(dim[0]) for j in range(dim[1])]

                shp_ctnt["ctnt"] = {"type": 'table', "dim": dim, "ctnt": cells}

            # SLide's images (no existing test functions)
            try:
                # Retrieve image
                img = shape.image
                blob, ext = img.blob, img.ext

                f_path = f'Outputs/{img_id}.{ext}'
                o_path = 'Outputs/out_' + str(img_id) + '.' + img.ext

                with open(f_path, 'wb') as file:
                    file.write(blob)

                # Sanitize file document
                imgs.ext_img(f_path, o_path)

                shp_ctnt["ctnt"] = {"type": 'img', "id": img_id, "file": o_path}

                img_id += 1

            except:
                pass

            s_ctnt["shapes"].append(shp_ctnt)

        layout.append(s_ctnt)

    return layout


################################################################################


def recompose(pres, layout):

    for slide in layout:

        # Add slide to pres
        sld_l = pres.slide_layouts[1]
        sld = pres.slides.add_slide(sld_l)
        left = top = Inches(1)

        width = height = Inches(5)

        # Add textbox to slide
        textBox = sld.shapes.add_textbox(left, top, width, height)
        textFrame = textBox.text_frame

        for shape in slide["shapes"]:

            if shape["ctnt"] != {}:

                if shape["ctnt"]["type"] == "text":

                    # Add text paragraph
                    p = textFrame.add_paragraph()
                    p.text = shape["ctnt"]["ctnt"]

                if shape["ctnt"]["type"] == "table":

                    # Get table dimensions
                    dim = shape["ctnt"]["dim"]

                    # Add table shape
                    t = sld.shapes.add_table(dim[0], dim[1], left, top, width, height)

                    # Fill table cells
                    for cell in shape["ctnt"]["ctnt"]:

                        c = t.table.cell(cell[0], cell[1])
                        c.text = cell[2]

                if shape["ctnt"]["type"] == "img":

                    # Add image
                    sld.shapes.add_picture(shape["ctnt"]["file"], Inches(1), Inches(1), Inches(5), Inches(5))

                    # Remove image file
                    remove(shape["ctnt"]["file"])

            else:
                pass


################################################################################


def sanitiz(path):

    info('Sanitizing ' + path.split("/")[-1])

    # Fetch layout from source
    try:
        l = get_layout(path)

    except:
        fail('Data extraction failed')
        return False, ''


    info('Creating new document')

    # Create sanitized document
    try:

        # Create presentation
        pres = Presentation()
        # Recompose document
        recompose(pres, l)
        # Save document
        pres.save('Outputs/out_' + path.split("/")[-1].split(".")[0] + ".pptx")

    except:
        fail('Document recomposition failed')
        return False, ''

    success('Document sanitized successfully.')
    return True, 'Outputs/out_' + path.split("/")[-1].split(".")[0] + ".pptx"


################################################################################


if __name__ == "__main__":
    print('Please run main.py or read software documentation')

    exit()

else:
    from .prints import *
    from . import imgs
